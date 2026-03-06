#!/usr/bin/env python3
"""
Link validation utility for Office documents.
Verifies that all hyperlinks in presentations and documents are accessible.
"""
import sys
import urllib.request
from pathlib import Path
from datetime import datetime

def extract_links_from_pptx(pptx_path):
    from pptx import Presentation
    
    links = []
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "click_action") and shape.click_action.hyperlink.address:
                    links.append(shape.click_action.hyperlink.address)
                # Check text frames for hyperlinks
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if hasattr(run, "hyperlink") and run.hyperlink.address:
                                links.append(run.hyperlink.address)
    except Exception as e:
        print(f"Error extracting links from PPTX: {e}")
    
    return links

def extract_links_from_docx(docx_path):
    from docx import Document
    
    links = []
    try:
        doc = Document(docx_path)
        
        # Extract from main document body
        for paragraph in doc.paragraphs:
            for run in paragraph.runs:
                if hasattr(run, "element") and run.element.xml:
                    # Check for hyperlinks in the XML
                    import re
                    hyperlinks = re.findall(r'<w:hyperlink[^>]*r:id="([^"]*)"', run.element.xml)
                    for rel_id in hyperlinks:
                        try:
                            rel = doc.part.rels[rel_id]
                            if rel.target_ref:
                                links.append(rel.target_ref)
                        except:
                            pass
        
        # Also check relationships directly
        for rel in doc.part.rels.values():
            if rel.reltype == "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink":
                if rel.target_ref:
                    links.append(rel.target_ref)
                    
    except Exception as e:
        print(f"Error extracting links from DOCX: {e}")
    
    return links

def extract_links(file_path):
    file_path = Path(file_path)
    suffix = file_path.suffix.lower()
    
    if suffix == '.pptx':
        return extract_links_from_pptx(file_path)
    elif suffix == '.docx':
        return extract_links_from_docx(file_path)
    else:
        print(f"Unsupported file type: {suffix}")
        return []

def main():
    # Create execution log immediately
    log_file = Path(".script_executed_link_tester.log")
    timestamp = datetime.now().isoformat()
    log_file.write_text(f"[{timestamp}] link_tester.py executed\n")
    
    if len(sys.argv) < 2:
        print("Usage: python link_tester.py <file>")
        print("Supports: .pptx, .docx")
        log_file.write_text(log_file.read_text() + f"[{timestamp}] ERROR: No file argument provided\n")
        sys.exit(1)
    
    file_path = sys.argv[1]
    log_file.write_text(log_file.read_text() + f"[{timestamp}] Target file: {file_path}\n")
    print(f"Testing links from: {file_path}")
    
    links = extract_links(file_path)
    
    if not links:
        print("No links found in file.")
        
        # Create validation report
        report_file = Path("link_validation_report.txt")
        report_file.write_text(f"Link validation completed for: {file_path}\nLinks tested: 0\n")
        return
    
    print(f"Found {len(links)} links. Testing them...")
    
    for idx, link in enumerate(links, 1):
        print(f"  {idx}. Testing: {link}")
        try:
            req = urllib.request.Request(link, headers={'User-Agent': 'Mozilla/5.0'})
            with urllib.request.urlopen(req, timeout=10) as response:
                status = response.getcode()
                print(f"     Status: {status} - {'OK' if status == 200 else 'Warning'}")
        except Exception as e:
            print(f"     Failed: {e}")
    
    print(f"\nAll {len(links)} links tested.")
    
    # Create validation report
    report_file = Path("link_validation_report.txt")
    report_file.write_text(f"Link validation completed for: {file_path}\nLinks tested: {len(links)}\n")

if __name__ == "__main__":
    main()
